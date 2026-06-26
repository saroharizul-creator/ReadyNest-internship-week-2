from fastapi import APIRouter, Depends, HTTPException, status, Query, Request
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
import pandas as pd
import io
import os
from datetime import datetime
from typing import Optional

# ReportLab imports for PDF
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# PPTX imports
from pptx import Presentation
from pptx.util import Inches, Pt

from openpyxl import Workbook

from app.database.connection import get_db
from app.database.models import Project, User
from app.auth.utils import decode_token
from app.analytics.routes import load_project_df
from app.analytics.engine import AnalyticsEngine
from app.analytics.eda_engine import EdaEngine
from app.analytics.product_engine import ProductEngine
from app.analytics.customer_engine import CustomerEngine
from app.analytics.regional_engine import RegionalEngine
from app.analytics.executive_summary_engine import ExecutiveSummaryEngine
from app.ml.engine import MLEngine

router = APIRouter(prefix="/projects", tags=["Reports Generation"])

# Helper dependency to authenticate either from header or query param
def get_current_user_from_header_or_query(
    request: Request,
    token: Optional[str] = Query(None),
    db: Session = Depends(get_db)
) -> User:
    token_str = None
    auth_header = request.headers.get("Authorization")
    if auth_header and auth_header.startswith("Bearer "):
        token_str = auth_header.split(" ")[1]
    elif token:
        token_str = token
        
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    if not token_str:
        raise credentials_exception
        
    payload = decode_token(token_str)
    if payload is None:
        raise credentials_exception
        
    email: str = payload.get("sub")
    if email is None:
        raise credentials_exception
        
    user = db.query(User).filter(User.email == email).first()
    if user is None:
        raise credentials_exception
        
    return user




@router.get("/{project_id}/reports/excel")
def generate_excel_report(
    project_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user_from_header_or_query)
):
    project = db.query(Project).filter(Project.id == project_id, Project.user_id == current_user.id).first()
    if not project:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Project not found")
        
    df = load_project_df(project_id, db)
    if df.empty:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No data available to export.")
        
    # Calculate all metrics
    bi_metrics = AnalyticsEngine.calculate_bi_metrics(df)
    eda_metrics = EdaEngine.run_profiling(df)
    prod_metrics = ProductEngine.run_analysis(df)
    cust_metrics = CustomerEngine.run_analysis(df)
    reg_metrics = RegionalEngine.run_analysis(df)
    ml_churn = MLEngine.predict_customer_churn_and_clv(df)
    ml_recs = MLEngine.generate_product_recommendations(df)
    
    # Create Excel buffer using optimized openpyxl write_only mode
    output = io.BytesIO()
    wb = Workbook(write_only=True)
    
    def append_df_to_wb(sheet_title, target_df):
        if target_df is None:
            return
        if isinstance(target_df, list):
            if not target_df:
                return
            target_df = pd.DataFrame(target_df)
        elif isinstance(target_df, pd.DataFrame) and target_df.empty:
            return
            
        ws = wb.create_sheet(title=sheet_title)
        ws.append(list(target_df.columns))
        clean_df = target_df.where(pd.notnull(target_df), None)
        for row in clean_df.itertuples(index=False, name=None):
            ws.append(row)
            
    # Sheet 1: Raw Cleaned Data
    append_df_to_wb("Cleaned Data", df)
    
    # Sheet 2: KPIs & Executive Insights
    kpis = bi_metrics.get("kpis", {})
    kpi_df = pd.DataFrame(list(kpis.items()), columns=["KPI Metric", "Value"])
    append_df_to_wb("Executive Summary", kpi_df)
    
    # Sheet 3: EDA Overview & Stats
    if eda_metrics:
        eda_overview = pd.DataFrame(list(eda_metrics["overview"].items()), columns=["Property", "Value"])
        append_df_to_wb("EDA Overview", eda_overview)
        
        stats_rows = []
        for col, summary in eda_metrics["stats_summary"].items():
            for stat_name, stat_val in summary.items():
                stats_rows.append({"Column": col, "Statistic": stat_name, "Value": stat_val})
        if stats_rows:
            append_df_to_wb("EDA Statistics", pd.DataFrame(stats_rows))
            
    # Sheet 4: Customer Performance
    if cust_metrics:
        cust_kpis = pd.DataFrame(list(cust_metrics["kpis"].items()), columns=["Customer Metric", "Value"])
        append_df_to_wb("Customer Analytics", cust_kpis)
        
        cust_timeline = pd.DataFrame(cust_metrics["acquisition_timeline"])
        append_df_to_wb("Acquisition Timeline", cust_timeline)
        
    # Sheet 5: Product Performance
    if prod_metrics:
        top_prod = pd.DataFrame(prod_metrics["top_selling"])
        append_df_to_wb("Product Analytics", top_prod)
        
        cat_anal = pd.DataFrame(prod_metrics["category_analysis"])
        append_df_to_wb("Category Analytics", cat_anal)
        
    # Sheet 6: Regional Performance
    if reg_metrics:
        reg_summ = pd.DataFrame(reg_metrics["region_summary"])
        append_df_to_wb("Regional Performance", reg_summ)
        
    # Sheet 7: Segmentation & ML Predicts
    rfm_df = pd.DataFrame(bi_metrics.get("rfm", []))
    append_df_to_wb("RFM Segments", rfm_df)
    
    append_df_to_wb("Churn & CLV Analysis", ml_churn)
    append_df_to_wb("Cross-Selling Models", ml_recs)
    
    wb.save(output)
    output.seek(0)
    
    filename = f"insightflow_enterprise_report_{project_id}_{datetime.now().strftime('%Y%m%d')}.xlsx"
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )


@router.get("/{project_id}/reports/pdf")
def generate_pdf_report(
    project_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user_from_header_or_query)
):
    project = db.query(Project).filter(Project.id == project_id, Project.user_id == current_user.id).first()
    if not project:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Project not found")
        
    df = load_project_df(project_id, db)
    if df.empty:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No data available to export.")
        
    # Calculate all required sections
    exec_summary = ExecutiveSummaryEngine.run_analysis(df)
    kpis = exec_summary.get("kpis", {})
    eda_metrics = EdaEngine.run_profiling(df)
    prod_metrics = ProductEngine.run_analysis(df)
    cust_metrics = CustomerEngine.run_analysis(df)
    reg_metrics = RegionalEngine.run_analysis(df)
    ml_churn = MLEngine.predict_customer_churn_and_clv(df)
    ml_recs = MLEngine.generate_product_recommendations(df)
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    story = []
    
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=22,
        leading=26,
        textColor=colors.HexColor('#0F172A'),
        spaceAfter=8
    )
    section_style = ParagraphStyle(
        'SecHeader',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=12,
        leading=16,
        textColor=colors.HexColor('#1E293B'),
        spaceBefore=14,
        spaceAfter=6
    )
    body_style = ParagraphStyle(
        'BodyTextCustom',
        parent=styles['BodyText'],
        fontName='Helvetica',
        fontSize=9,
        leading=12,
        textColor=colors.HexColor('#334155')
    )
    bold_body_style = ParagraphStyle(
        'BoldBodyCustom',
        parent=body_style,
        fontName='Helvetica-Bold'
    )
    
    # 1. Title Page
    story.append(Paragraph(f"InsightFlow Enterprise Analytics Report", title_style))
    story.append(Paragraph(f"Project Workgroup: {project.name} | Industry Sector: {project.industry or 'General'}", body_style))
    story.append(Paragraph(f"Report Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')} (IST)", body_style))
    story.append(Spacer(1, 10))
    
    # 2. Executive Summary KPIs Table
    story.append(Paragraph("I. Executive Key Performance Indicators", section_style))
    kpi_data = [
        [Paragraph("<b>Financial Metric</b>", body_style), Paragraph("<b>Value</b>", body_style)],
        [Paragraph("Total Net Sales", body_style), Paragraph(f"${kpis.get('total_revenue', 0.0):,.2f}", body_style)],
        [Paragraph("Total Invoice Orders", body_style), Paragraph(f"{kpis.get('total_orders', 0):,}", body_style)],
        [Paragraph("Active Customer Pool", body_style), Paragraph(f"{kpis.get('total_customers', 0):,}", body_style)],
        [Paragraph("Average Order Value (AOV)", body_style), Paragraph(f"${kpis.get('aov', 0.0):,.2f}", body_style)],
        [Paragraph("Average Gross Profit Margin", body_style), Paragraph(f"{kpis.get('profit_margin', 0.0)*100:.1f}%", body_style)],
        [Paragraph("Customer Retention Rate", body_style), Paragraph(f"{kpis.get('retention_rate', 0.0)*100:.1f}%", body_style)]
    ]
    t = Table(kpi_data, colWidths=[270, 270])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#F8FAFC')),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('BOTTOMPADDING', (0,0), (-1,0), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
        ('TOPPADDING', (0,1), (-1,-1), 4),
        ('BOTTOMPADDING', (0,1), (-1,-1), 4),
    ]))
    story.append(t)
    story.append(Spacer(1, 10))
    
    # 3. AI Strategic Insights
    story.append(Paragraph("II. Executive Strategic Insights & Recommendations", section_style))
    story.append(Paragraph("<b>Top Key Insights:</b>", bold_body_style))
    for ins in exec_summary.get("insights", [])[:3]:
        story.append(Paragraph(f"&bull; {ins}", body_style))
    story.append(Spacer(1, 5))
    story.append(Paragraph("<b>Top Recommendations:</b>", bold_body_style))
    for rec in exec_summary.get("recommendations", [])[:3]:
        story.append(Paragraph(f"&bull; {rec}", body_style))
    story.append(Spacer(1, 5))
    story.append(Paragraph("<b>Identified Risks & Opportunities:</b>", bold_body_style))
    for risk in exec_summary.get("risks", [])[:2]:
        story.append(Paragraph(f"&bull; <i>Risk:</i> {risk}", body_style))
    for opp in exec_summary.get("opportunities", [])[:2]:
        story.append(Paragraph(f"&bull; <i>Opportunity:</i> {opp}", body_style))
    
    story.append(PageBreak())

    # 4. EDA Findings
    if eda_metrics:
        story.append(Paragraph("III. Exploratory Data Analysis (EDA) Summary", section_style))
        overview = eda_metrics.get("overview", {})
        eda_data = [
            [Paragraph("<b>EDA Property</b>", body_style), Paragraph("<b>Value</b>", body_style)],
            [Paragraph("Total Row Count", body_style), Paragraph(f"{overview.get('row_count', 0):,}", body_style)],
            [Paragraph("Total Column Count", body_style), Paragraph(f"{overview.get('column_count', 0)}", body_style)],
            [Paragraph("Duplicate Records Count", body_style), Paragraph(f"{overview.get('duplicate_count', 0)}", body_style)],
            [Paragraph("Data Quality Integrity Score", body_style), Paragraph(f"{overview.get('quality_score', 0.0)}%", body_style)]
        ]
        et = Table(eda_data, colWidths=[270, 270])
        et.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#F8FAFC')),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
            ('TOPPADDING', (0,0), (-1,-1), 4),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ]))
        story.append(et)
        story.append(Spacer(1, 10))

    # 5. Product Performance
    if prod_metrics:
        story.append(Paragraph("V. Product Performance Breakdown", section_style))
        story.append(Paragraph("<b>Top Revenue Performing Products:</b>", bold_body_style))
        top_p_data = [[Paragraph("<b>Product Name</b>", body_style), Paragraph("<b>Category</b>", body_style), Paragraph("<b>Revenue</b>", body_style), Paragraph("<b>Contribution</b>", body_style)]]
        for p in prod_metrics.get("top_selling", [])[:5]:
            top_p_data.append([
                Paragraph(p["product_name"], body_style),
                Paragraph(p["category"], body_style),
                Paragraph(f"${p['revenue']:,.2f}", body_style),
                Paragraph(f"{p['contribution_pct']:.1f}%", body_style)
            ])
        pt = Table(top_p_data, colWidths=[220, 120, 100, 100])
        pt.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#F8FAFC')),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
            ('TOPPADDING', (0,0), (-1,-1), 4),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ]))
        story.append(pt)
        story.append(Spacer(1, 10))

    # 6. Customer Overview & Segmentation
    if cust_metrics:
        story.append(Paragraph("V. Customer Overview & Retention Trends", section_style))
        cust_kpis = cust_metrics.get("kpis", {})
        story.append(Paragraph(f"&bull; Total Customer Base Size: <b>{cust_kpis.get('total_customers', 0):,}</b> unique buyers.", body_style))
        story.append(Paragraph(f"&bull; Repeat Purchase Rate: <b>{cust_kpis.get('repeat_purchase_rate', 0.0):.1f}%</b>.", body_style))
        story.append(Spacer(1, 10))

    # 7. Regional Performance
    if reg_metrics:
        story.append(Paragraph("VI. Regional Performance Breakdown", section_style))
        reg_data = [[Paragraph("<b>Geographical Region</b>", body_style), Paragraph("<b>Gross Sales</b>", body_style), Paragraph("<b>Net Profit</b>", body_style), Paragraph("<b>Order Count</b>", body_style)]]
        for r in reg_metrics.get("region_summary", []):
            reg_data.append([
                Paragraph(r["region"], body_style),
                Paragraph(f"${r['revenue']:,.2f}", body_style),
                Paragraph(f"${r['profit']:,.2f}", body_style),
                Paragraph(f"{r['orders']:,}", body_style)
            ])
        rt = Table(reg_data, colWidths=[180, 120, 120, 120])
        rt.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#F8FAFC')),
            ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#CBD5E1')),
            ('TOPPADDING', (0,0), (-1,-1), 4),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ]))
        story.append(rt)
        story.append(Spacer(1, 10))

    # 8. ML Predictions: Churn & Recommendations
    story.append(Paragraph("VII. Predictive Intelligence & Cross-Selling Opportunity", section_style))
    if ml_recs:
        story.append(Paragraph("<b>Top Cross-Selling Opportunities (Bought X also bought Y):</b>", bold_body_style))
        for idx, rec in enumerate(ml_recs[:3]):
            story.append(Paragraph(f"&bull; Recommend <b>{rec['recommended']}</b> to buyers of <b>{rec['product']}</b> (Correlation Support: <b>{rec['confidence']*100:.0f}%</b>).", body_style))
    story.append(Spacer(1, 10))
    if ml_churn:
        high_risk_count = sum(1 for c in ml_churn if c["riskLevel"] == "High")
        story.append(Paragraph(f"&bull; Churn Risk Assessment: Detected <b>{high_risk_count}</b> customers in the **High Churn Risk** category.", body_style))

    doc.build(story)
    buffer.seek(0)
    
    filename = f"insightflow_enterprise_report_{project_id}_{datetime.now().strftime('%Y%m%d')}.pdf"
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )


@router.get("/{project_id}/reports/powerpoint")
def generate_pptx_report(
    project_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user_from_header_or_query)
):
    project = db.query(Project).filter(Project.id == project_id, Project.user_id == current_user.id).first()
    if not project:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Project not found")
        
    df = load_project_df(project_id, db)
    if df.empty:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No data available to export.")
        
    # Gather analytics summary
    exec_summary = ExecutiveSummaryEngine.run_analysis(df)
    kpis = exec_summary.get("kpis", {})
    
    prs = Presentation()
    
    # Slide 1: Cover Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"InsightFlow Executive Presentation"
    subtitle.text = f"Workspace: {project.name}\nGenerated on {datetime.now().strftime('%Y-%m-%d')}\nConfidential Executive Report"
    
    # Slide 2: KPI Dashboard summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Key Performance Indicators Summary"
    tf = shapes.placeholders[1].text_frame
    tf.text = f"Total Sales Revenue: ${kpis.get('total_revenue', 0.0):,.2f}"
    tf.add_paragraph().text = f"Total Orders Placed: {kpis.get('total_orders', 0):,}"
    tf.add_paragraph().text = f"Active Customer Pool: {kpis.get('total_customers', 0):,}"
    tf.add_paragraph().text = f"Average Order Value (AOV): ${kpis.get('aov', 0.0):,.2f}"
    tf.add_paragraph().text = f"Operating Margin efficiency: {kpis.get('profit_margin', 0.0)*100:.1f}%"
    tf.add_paragraph().text = f"Repeat Purchase Retention Rate: {kpis.get('retention_rate', 0.0)*100:.1f}%"

    # Slide 3: Executive Strategic Insights
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Strategic Insights & Business Opportunities"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Key Strategic Findings:"
    for ins in exec_summary.get("insights", [])[:2]:
        p = tf.add_paragraph()
        p.text = f"- {ins}"
        p.level = 1
    p_opp = tf.add_paragraph()
    p_opp.text = "Key Opportunities Identified:"
    p_opp.space_before = Pt(12)
    for opp in exec_summary.get("opportunities", [])[:2]:
        p = tf.add_paragraph()
        p.text = f"- {opp}"
        p.level = 1

    # Slide 4: Recommendations & Growth Plans
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    shapes.title.text = "Actionable Recommendations"
    tf = shapes.placeholders[1].text_frame
    tf.text = "Growth Initiatives:"
    for rec in exec_summary.get("recommendations", [])[:3]:
        p = tf.add_paragraph()
        p.text = f"- {rec}"
        p.level = 1
    p_rev = tf.add_paragraph()
    p_rev.text = "Revenue Suggestions:"
    p_rev.space_before = Pt(12)
    for sug in exec_summary.get("revenue_suggestions", [])[:2]:
        p = tf.add_paragraph()
        p.text = f"- {sug}"
        p.level = 1

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    filename = f"insightflow_presentation_{project_id}_{datetime.now().strftime('%Y%m%d')}.pptx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )
